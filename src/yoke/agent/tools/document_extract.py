"""Tool for extracting readable context from various document formats."""

from __future__ import annotations

import html
import importlib
import re
import zipfile
from pathlib import Path

from pydantic import Field

from yoke.agent.tools.base import WorkspaceTool
from yoke.agent.truncate import truncate_head

TEXT_EXTENSIONS = {
    ".csv",
    ".css",
    ".html",
    ".ini",
    ".js",
    ".json",
    ".log",
    ".md",
    ".py",
    ".rst",
    ".scss",
    ".toml",
    ".ts",
    ".tsx",
    ".txt",
    ".xml",
    ".yaml",
    ".yml",
}


def _read_utf8_best_effort(path: Path) -> str:
    raw = path.read_bytes()
    for encoding in (
        "utf-8-sig",
        "utf-16",
        "utf-16-le",
        "utf-16-be",
        "utf-8",
        "cp1252",
    ):
        try:
            return raw.decode(encoding).replace("\r\n", "\n").replace("\r", "\n")
        except UnicodeDecodeError:
            continue
    return (
        raw.decode("utf-8", errors="replace").replace("\r\n", "\n").replace("\r", "\n")
    )


def _is_probably_binary(path: Path) -> bool:
    sample = path.read_bytes()[:4096]
    if not sample:
        return False
    if b"\x00" in sample:
        return True
    textish = sum(1 for byte in sample if byte in b"\t\n\r\f\b" or 32 <= byte <= 126)
    return (textish / len(sample)) < 0.75


def _extract_docx(path: Path) -> str:
    try:
        document_module = importlib.import_module("docx")
    except ImportError:
        return _extract_office_xml(
            path,
            member_names=("document.xml", "header", "footer"),
            paragraph_tags=("w:p",),
        )

    try:
        doc = document_module.Document(path)
        lines: list[str] = []
        for paragraph in doc.paragraphs:
            if paragraph.text.strip():
                lines.append(paragraph.text)
        for table in doc.tables:
            for row in table.rows:
                row_text = " | ".join(cell.text.strip() for cell in row.cells)
                if row_text.strip():
                    lines.append(row_text)
        for section in doc.sections:
            for paragraph in section.header.paragraphs:
                if paragraph.text.strip():
                    lines.append(f"[HEADER] {paragraph.text}")
            for paragraph in section.footer.paragraphs:
                if paragraph.text.strip():
                    lines.append(f"[FOOTER] {paragraph.text}")
        return "\n".join(lines) if lines else f"[No text detected in {path.name}]"
    except Exception as exc:
        return f"[Error extracting {path.name}: {str(exc)[:160]}]"


def _extract_xlsx(path: Path) -> str:
    try:
        from openpyxl import load_workbook  # type: ignore[import-untyped]  # ty: ignore[unresolved-import]
    except ImportError:
        return _extract_xlsx_xml(path)

    try:
        workbook = load_workbook(path, data_only=True)
        lines: list[str] = []
        for sheet_name in workbook.sheetnames:
            worksheet = workbook[sheet_name]
            lines.append(f"Sheet: {sheet_name}")
            for row in worksheet.iter_rows(values_only=True):
                row_text = " | ".join(
                    "" if value is None else str(value) for value in row
                )
                if row_text.strip():
                    lines.append(row_text)
        return "\n".join(lines) if lines else f"[No text detected in {path.name}]"
    except Exception as exc:
        return f"[Error extracting {path.name}: {str(exc)[:160]}]"


def _extract_pdf(path: Path) -> str:
    try:
        import pdfplumber
    except ImportError:
        return (
            "[Missing dependency: pdfplumber. Install with `pip install pdfplumber`.]"
        )

    try:
        lines: list[str] = []
        with pdfplumber.open(path) as pdf:
            for index, page in enumerate(pdf.pages, start=1):
                text = page.extract_text()
                if text:
                    lines.append(f"Page {index}:\n{text}")
        return "\n\n".join(lines) if lines else f"[No text detected in {path.name}]"
    except Exception as exc:
        return f"[Error extracting {path.name}: {str(exc)[:160]}]"


def _extract_pptx(path: Path) -> str:
    try:
        from pptx import Presentation  # ty: ignore[unresolved-import]
    except ImportError:
        return _extract_office_xml(
            path,
            member_names=("slide",),
            paragraph_tags=("a:p",),
        )

    try:
        presentation = Presentation(str(path))
        lines: list[str] = []
        for index, slide in enumerate(presentation.slides, start=1):
            lines.append(f"Slide {index}:")
            for shape in slide.shapes:
                text = getattr(shape, "text", "")
                if isinstance(text, str) and text.strip():
                    lines.append(text)
        return "\n".join(lines) if lines else f"[No text detected in {path.name}]"
    except Exception as exc:
        return f"[Error extracting {path.name}: {str(exc)[:160]}]"


def _extract_image(path: Path) -> str:
    try:
        easyocr_module = importlib.import_module("easyocr")
    except ImportError:
        return "[Missing dependency: easyocr. Install with `pip install easyocr`.]"

    try:
        reader = easyocr_module.Reader(["en"], gpu=False)
        results = reader.readtext(str(path))
        text = "\n".join(item[1] for item in results if len(item) > 1)
        return text if text.strip() else f"[No text detected in {path.name}]"
    except Exception as exc:
        return f"[Error extracting {path.name}: {str(exc)[:160]}]"


def _extract_office_xml(
    path: Path,
    *,
    member_names: tuple[str, ...],
    paragraph_tags: tuple[str, ...],
) -> str:
    try:
        with zipfile.ZipFile(path) as archive:
            lines: list[str] = []
            for member_name in archive.namelist():
                lowered = member_name.lower()
                if not any(token in lowered for token in member_names):
                    continue
                raw = archive.read(member_name).decode("utf-8", errors="replace")
                text = raw
                for tag in paragraph_tags:
                    text = text.replace(f"</{tag}>", "\n")
                text = html.unescape(
                    "".join(
                        ch if ch.isprintable() or ch == "\n" else " " for ch in text
                    )
                )
                compact = re.sub(r"<[^>]+>", " ", text)
                compact = re.sub(r"[ \t]+", " ", compact)
                compact = re.sub(r"\n{2,}", "\n", compact).strip()
                if compact:
                    lines.append(compact)
        return "\n".join(lines) if lines else f"[No text detected in {path.name}]"
    except Exception as exc:
        return f"[Error extracting {path.name}: {str(exc)[:160]}]"


def _extract_xlsx_xml(path: Path) -> str:
    try:
        with zipfile.ZipFile(path) as archive:
            shared_strings: list[str] = []
            if "xl/sharedStrings.xml" in archive.namelist():
                shared_raw = archive.read("xl/sharedStrings.xml").decode(
                    "utf-8", errors="replace"
                )
                shared_strings = [
                    html.unescape(re.sub(r"<[^>]+>", " ", match)).strip()
                    for match in re.findall(
                        r"<t[^>]*>(.*?)</t>", shared_raw, flags=re.DOTALL
                    )
                ]

            lines: list[str] = []
            for member_name in sorted(archive.namelist()):
                if not member_name.startswith(
                    "xl/worksheets/"
                ) or not member_name.endswith(".xml"):
                    continue
                sheet_raw = archive.read(member_name).decode("utf-8", errors="replace")
                lines.append(f"Sheet: {Path(member_name).stem}")
                for row_match in re.findall(
                    r"<row[^>]*>(.*?)</row>", sheet_raw, flags=re.DOTALL
                ):
                    values: list[str] = []
                    for cell_match in re.findall(
                        r"<c([^>]*)>(.*?)</c>", row_match, flags=re.DOTALL
                    ):
                        attrs, body = cell_match
                        value_match = re.search(r"<v>(.*?)</v>", body, flags=re.DOTALL)
                        if value_match is None:
                            continue
                        value = html.unescape(value_match.group(1)).strip()
                        if ' t="s"' in attrs and value.isdigit():
                            index = int(value)
                            value = (
                                shared_strings[index]
                                if 0 <= index < len(shared_strings)
                                else value
                            )
                        values.append(value)
                    row_text = " | ".join(value for value in values if value)
                    if row_text:
                        lines.append(row_text)
        return "\n".join(lines) if lines else f"[No text detected in {path.name}]"
    except Exception as exc:
        return f"[Error extracting {path.name}: {str(exc)[:160]}]"


EXTRACTORS = {
    ".docx": _extract_docx,
    ".xlsx": _extract_xlsx,
    ".pdf": _extract_pdf,
    ".pptx": _extract_pptx,
    ".png": _extract_image,
    ".jpg": _extract_image,
    ".jpeg": _extract_image,
    ".gif": _extract_image,
    ".bmp": _extract_image,
    ".webp": _extract_image,
    ".tiff": _extract_image,
    ".tif": _extract_image,
}


class ExtractFileContextTool(WorkspaceTool):
    """Extract readable context from a file using the best strategy."""

    name = "extract_file_context"
    description = (
        "Extract readable context from a file under the workspace root. "
        "Supports plain text directly and best-effort extraction for "
        "docx, xlsx, pdf, pptx, and common image formats."
    )

    path: str = Field(min_length=1)
    max_chars: int = Field(default=20_000, ge=500, le=200_000)

    def execute(self) -> dict[str, object]:
        """Extract content from the file and return it."""
        try:
            path = self._resolve_path(self.path)
            if not path.is_file():
                raise ValueError(f"Path is not a regular file: {self.path}")
            suffix = path.suffix.lower()
            if suffix in TEXT_EXTENSIONS:
                extracted = _read_utf8_best_effort(path)
                extractor = "text"
            elif suffix in EXTRACTORS:
                extracted = EXTRACTORS[suffix](path)
                extractor = suffix.lstrip(".")
            elif suffix in {".docm", ".dotx", ".dotm"}:
                extracted = _extract_office_xml(
                    path,
                    member_names=("document.xml", "header", "footer"),
                    paragraph_tags=("w:p",),
                )
                extractor = "office_xml"
            else:
                if _is_probably_binary(path):
                    return self._success(
                        extractor="binary",
                        content=f"[Unsupported binary file: {path.name}]",
                        details={"sizeBytes": path.stat().st_size},
                    )
                extracted = _read_utf8_best_effort(path)
                extractor = "best_effort_text"

            limited = (
                extracted[: self.max_chars]
                if len(extracted) > self.max_chars
                else extracted
            )
            truncation = truncate_head(limited)
            content = truncation.content
            truncated = truncation.truncated or len(extracted) > self.max_chars
            if len(extracted) > self.max_chars and not truncation.truncated:
                content = limited.rstrip() + "\n\n[Output truncated by max_chars.]"
            result = self._success(
                extractor=extractor,
                content=content,
                details={"sizeBytes": path.stat().st_size, "suffix": suffix},
            )
            if truncated:
                result["truncated"] = True
            return result
        except Exception as exc:
            return self._error(str(exc), path=self.path)
